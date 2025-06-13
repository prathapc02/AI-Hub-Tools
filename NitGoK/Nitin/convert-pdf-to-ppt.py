import sys

try:
    from tkinter import *
    from tkinter import filedialog, messagebox
    import threading
    from PyPDF2 import PdfReader
    from pptx import Presentation
    from pptx.util import Inches
except ModuleNotFoundError as error:
    print(error)
    input("Press Any Key To Exit .. ")
    sys.exit()


class PdfToPpt:
    def __init__(self):
        # Window Design
        self.root = Tk()
        self.root.title("PDF2PPT")
        self.root.geometry("400x300")
        self.root.resizable(0, 0)
        self.root.configure(bg="#FFF")
        self.root.iconbitmap("ppt-icon.ico")

        # Vars
        self.pdf_var = StringVar()
        self.ppt_var = StringVar()

        # Title Label
        title_label = Label(self.root, text="Convert PDF to PPT", bg="#FFF", font=("roboto", 20, "bold"))
        title_label.pack(fill="x", pady=20)

        # Entries
        pdf_entry = Entry(self.root, width=35, highlightbackground="black", state="disabled",
                          highlightcolor="black", highlightthickness=2, textvariable=self.pdf_var)
        pdf_entry.place(x=30, y=120)

        ppt_entry = Entry(self.root, width=35, highlightbackground="black", state="disabled",
                          highlightcolor="black", highlightthickness=2, textvariable=self.ppt_var)
        ppt_entry.place(x=30, y=160)

        # Buttons
        pdf_button = Button(self.root, text="PDF", width=10, borderwidth=2, bg="black",
                            fg="white", font=("verdana", 10, "bold"), command=self.get_pdf)
        pdf_button.place(x=265, y=118)

        save_button = Button(self.root, text="Save", width=10, borderwidth=2, bg="black",
                             fg="white", font=("verdana", 10, "bold"), command=self.save)
        save_button.place(x=265, y=158)

        convert_button = Button(self.root, text="Convert", width=10, borderwidth=2, bg="black",
                                fg="white", font=("verdana", 10, "bold"), command=self.convert_thread)
        convert_button.place(x=150, y=220)

        clear_button = Button(self.root, text="Clear", width=10, borderwidth=2, bg="black",
                              fg="white", font=("verdana", 10, "bold"), command=self.clear)
        clear_button.place(x=265, y=220)

        exit_button = Button(self.root, text="Exit", width=10, borderwidth=2, bg="black",
                             fg="white", font=("verdana", 10, "bold"), command=self.exit_)
        exit_button.place(x=30, y=220)

        self.root.mainloop()

    def get_pdf(self):
        self.pdf_var.set(filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf")]))

    def save(self):
        self.ppt_var.set(filedialog.asksaveasfilename(defaultextension=".pptx",
                                                      filetypes=[("PPTX files", "*.pptx")]))

    def convert(self):
        if not self.pdf_var.get():
            messagebox.showerror("Error", "Choose PDF File")
        elif not self.ppt_var.get():
            messagebox.showerror("Error", "Choose Save Path")
        else:
            try:
                reader = PdfReader(self.pdf_var.get())
                presentation = Presentation()

                for page in reader.pages:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                    text = page.extract_text()

                    if text:
                        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                        text_frame = text_box.text_frame
                        text_frame.text = text
                    else:
                        messagebox.showwarning("Warning", "One or more pages are empty and were skipped.")

                presentation.save(self.ppt_var.get())
                messagebox.showinfo("Converted", "PDF successfully converted to PPT.")
            except Exception as e:
                print(e)
                messagebox.showerror("Error", "An error occurred during conversion.")

    def convert_thread(self):
        threading.Thread(target=self.convert).start()

    def clear(self):
        self.pdf_var.set("")
        self.ppt_var.set("")

    def exit_(self):
        check = messagebox.askyesno("Exit", "Do you want to exit?")
        if check:
            self.root.destroy()


if __name__ == '__main__':
    PdfToPpt()
