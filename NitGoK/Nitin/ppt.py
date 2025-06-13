import tkinter as tk
from tkinter import simpledialog, messagebox, filedialog
import wikipedia
from pptx import Presentation

def fetch_wikipedia_content(topic):
    """
    Fetch content from Wikipedia for the given topic.
    """
    try:
        # Fetch Wikipedia summary
        raw_content = wikipedia.summary(topic, sentences=5)  # Limit to 5 sentences
        return raw_content
    except wikipedia.exceptions.PageError:
        return "No Wikipedia page found for this topic."
    except wikipedia.exceptions.DisambiguationError as e:
        return f"Topic is ambiguous. Suggestions: {', '.join(e.options[:5])}"

def generate_powerpoint_from_agenda(agenda_points):
    """
    Generate a PowerPoint presentation based on the given agenda.
    """
    try:
        # Ask user to save the PowerPoint file
        file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])
        if not file_path:
            return

        # Create PowerPoint presentation
        presentation = Presentation()

        # Add Title Slide
        slide_layout = presentation.slide_layouts[0]  # Title Slide Layout
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Presentation Based on Agenda"
        subtitle.text = "Generated with Wikipedia Content"

        # Add Agenda Slide
        agenda_slide_layout = presentation.slide_layouts[1]  # Title + Content Layout
        agenda_slide = presentation.slides.add_slide(agenda_slide_layout)
        agenda_title = agenda_slide.shapes.title
        agenda_content = agenda_slide.placeholders[1]
        agenda_title.text = "Agenda"
        agenda_content.text = "\n".join([f"{i+1}. {point}" for i, point in enumerate(agenda_points)])

        # Generate slides for each agenda point
        for i, point in enumerate(agenda_points):
            # Fetch content for the agenda point
            content = fetch_wikipedia_content(point)
            if not content or "No Wikipedia page" in content:
                content = "Content not available for this topic."

            # Add a new slide for the agenda point
            slide = presentation.slides.add_slide(agenda_slide_layout)
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]
            title.text = f"Topic {i+1}: {point}"
            content_placeholder.text = content

        # Save PowerPoint file
        presentation.save(file_path)
        messagebox.showinfo("Success", f"PowerPoint file saved successfully:\n{file_path}")

    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {str(e)}")

def input_agenda_for_ppt():
    """
    Open a dialog box to get agenda points and generate PowerPoint content.
    """
    agenda_points = []
    while True:
        point = simpledialog.askstring("Input Agenda", "Enter an agenda point (or press Cancel to finish):")
        if not point:
            break
        agenda_points.append(point)

    if not agenda_points:
        messagebox.showwarning("No Input", "Please enter at least one agenda point.")
        return

    generate_powerpoint_from_agenda(agenda_points)

# Create the main window
root = tk.Tk()
root.title("Agenda-Based PowerPoint Generator")
root.geometry("400x200")

# Label for the app title
title_label = tk.Label(root, text="Agenda-Based PowerPoint Generator", font=("Helvetica", 16), wraplength=380)
title_label.pack(pady=20)

# Button to trigger PowerPoint generation
ppt_button = tk.Button(root, text="Generate PowerPoint", width=20, height=2, bg="white", fg="#262626", command=input_agenda_for_ppt)
ppt_button.pack(pady=20)

# Run the main loop
root.mainloop()
